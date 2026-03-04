#!/usr/bin/env python3
"""
Are.na Channel → Google Slides pipeline.

Paste an Are.na channel URL, and this script will:
  1. Download all images from the channel
  2. Build a 16:9 PowerPoint deck (dark gray background, images scaled to fit)
  3. Upload it to Google Drive as a Google Slides presentation

Usage:
    python arena_to_slides.py
    python arena_to_slides.py https://www.are.na/james-hicks/baroque-user-interface
    python arena_to_slides.py baroque-user-interface

Requirements (install in venv):
    pip install requests python-pptx Pillow

External tool:
    brew install rclone   (must be configured with: rclone config create gdrive drive)
"""

import os
import re
import sys
import json
import time
import shutil
import subprocess
from urllib.parse import urlparse

import requests
from pptx import Presentation
from pptx.util import Inches
from pptx.dml.color import RGBColor
from PIL import Image

# ── constants ────────────────────────────────────────────────────────────────

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG_COLOR = RGBColor(0x33, 0x33, 0x33)
SUPPORTED_EXT = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"}
BASE_DIR = os.path.dirname(os.path.abspath(__file__))

# ── helpers ──────────────────────────────────────────────────────────────────

def get_slug(input_str):
    """Extract channel slug from a URL or bare slug string."""
    input_str = input_str.strip()
    if "are.na" in input_str:
        return input_str.rstrip("/").split("/")[-1]
    return input_str


def prompt_for_url():
    """Ask the user to paste a channel URL interactively."""
    print("╭──────────────────────────────────────────╮")
    print("│  Are.na → Google Slides                  │")
    print("╰──────────────────────────────────────────╯")
    print()
    url = input("Paste an Are.na channel URL (or slug): ").strip()
    if not url:
        print("No input provided. Exiting.")
        sys.exit(1)
    return url

# ── step 1: download images ─────────────────────────────────────────────────

def download_channel_images(slug, output_dir):
    """Download all Image blocks from an Are.na channel. Returns image count."""
    os.makedirs(output_dir, exist_ok=True)
    base_url = f"https://api.are.na/v2/channels/{slug}/contents"
    page = 1
    per_page = 50
    downloaded = 0
    failed = 0
    total_blocks = 0

    print(f"\n[1/3] Downloading images from channel: {slug}\n")

    while True:
        resp = requests.get(base_url, params={"page": page, "per": per_page})
        if resp.status_code != 200:
            print(f"  Error fetching page {page}: {resp.status_code}")
            break

        data = resp.json()
        contents = data.get("contents", [])
        if not contents:
            break

        if page == 1:
            print(f"  Total blocks in channel: {data.get('length', '?')}")

        for block in contents:
            total_blocks += 1
            if block.get("class") != "Image":
                continue

            block_id = block.get("id", "unknown")
            title = block.get("title") or ""
            image = block.get("image", {})

            image_url = (
                image.get("original", {}).get("url")
                or image.get("large", {}).get("url")
                or image.get("display", {}).get("url")
                or image.get("thumb", {}).get("url")
            )
            if not image_url and isinstance(block.get("source"), dict):
                image_url = block["source"].get("url")
            if not image_url:
                continue

            if image_url.startswith("//"):
                image_url = "https:" + image_url
            elif not image_url.startswith("http"):
                image_url = "https://images.are.na/" + image_url.lstrip("/")

            ext = os.path.splitext(urlparse(image_url).path)[1] or ".jpg"
            safe_title = re.sub(r'[^\w\s-]', '', title)[:60].strip() if title else ""
            filename = f"{block_id}_{safe_title}{ext}" if safe_title else f"{block_id}{ext}"
            filename = filename.replace(" ", "_")
            filepath = os.path.join(output_dir, filename)

            if os.path.exists(filepath) and os.path.getsize(filepath) > 0:
                downloaded += 1
                continue

            try:
                img_resp = requests.get(image_url, timeout=30, headers={
                    "User-Agent": "Mozilla/5.0 (Macintosh; Intel Mac OS X 10_15_7)"
                })
                img_resp.raise_for_status()
                if len(img_resp.content) == 0:
                    failed += 1
                    continue
                with open(filepath, "wb") as f:
                    f.write(img_resp.content)
                downloaded += 1
                size_kb = len(img_resp.content) / 1024
                print(f"  Downloaded: {filename} ({size_kb:.1f} KB)")
            except Exception as e:
                print(f"  Failed: {block_id} – {e}")
                failed += 1

            time.sleep(0.2)

        page += 1

    print(f"  Done: {downloaded} images downloaded, {failed} failed")
    return downloaded

# ── step 2: build slideshow ──────────────────────────────────────────────────

def build_slideshow(image_dir, output_path):
    """Create a 16:9 .pptx with one image per slide on a dark gray background."""
    print(f"\n[2/3] Building slideshow\n")

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    blank_layout = prs.slide_layouts[6]

    files = sorted(os.listdir(image_dir))
    images = [f for f in files if os.path.splitext(f)[1].lower() in SUPPORTED_EXT]

    if not images:
        print("  No images found – nothing to do.")
        return None

    count = 0
    for fname in images:
        fpath = os.path.join(image_dir, fname)
        try:
            with Image.open(fpath) as img:
                img_w, img_h = img.size
        except Exception:
            continue

        slide = prs.slides.add_slide(blank_layout)
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        img_aspect = img_w / img_h
        slide_aspect = SLIDE_W / SLIDE_H

        if img_aspect > slide_aspect:
            pic_w = SLIDE_W
            pic_h = int(SLIDE_W / img_aspect)
        else:
            pic_h = SLIDE_H
            pic_w = int(SLIDE_H * img_aspect)

        left = (SLIDE_W - pic_w) // 2
        top = (SLIDE_H - pic_h) // 2
        slide.shapes.add_picture(fpath, left, top, pic_w, pic_h)
        count += 1

    prs.save(output_path)
    print(f"  Created {count} slides → {os.path.basename(output_path)}")
    return output_path

# ── step 3: upload to google drive ───────────────────────────────────────────

def upload_to_google_drive(pptx_path):
    """Upload .pptx to Google Drive via rclone, converting to Google Slides."""
    print(f"\n[3/3] Uploading to Google Drive\n")

    if not shutil.which("rclone"):
        print("  Error: rclone is not installed. Install with: brew install rclone")
        print(f"  Your .pptx is saved locally at: {pptx_path}")
        return False

    # Check that the gdrive remote exists
    result = subprocess.run(["rclone", "listremotes"], capture_output=True, text=True)
    if "gdrive:" not in result.stdout:
        print("  Error: rclone remote 'gdrive' not configured.")
        print("  Run: rclone config create gdrive drive")
        print(f"  Your .pptx is saved locally at: {pptx_path}")
        return False

    cmd = [
        "rclone", "copy", pptx_path, "gdrive:",
        "--drive-import-formats", "pptx",
        "-v"
    ]
    proc = subprocess.run(cmd, capture_output=True, text=True)

    if proc.returncode == 0:
        name = os.path.splitext(os.path.basename(pptx_path))[0]
        print(f"  Uploaded and converted to Google Slides: \"{name}\"")
        print(f"  Open Google Drive to view it.")
        return True
    else:
        print(f"  Upload failed: {proc.stderr}")
        print(f"  Your .pptx is saved locally at: {pptx_path}")
        return False

# ── main ─────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) > 1:
        raw = sys.argv[1]
    else:
        raw = prompt_for_url()

    slug = get_slug(raw)
    image_dir = os.path.join(BASE_DIR, f"arena_{slug}")
    pptx_path = os.path.join(BASE_DIR, f"{slug}-slides.pptx")

    # Step 1 – download
    count = download_channel_images(slug, image_dir)
    if count == 0:
        print("\nNo images downloaded. Exiting.")
        sys.exit(1)

    # Step 2 – build slideshow
    result = build_slideshow(image_dir, pptx_path)
    if not result:
        sys.exit(1)

    # Step 3 – upload
    upload_to_google_drive(pptx_path)

    print("\nAll done!")


if __name__ == "__main__":
    main()
