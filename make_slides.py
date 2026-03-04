#!/usr/bin/env python3
"""Create a 16:9 PowerPoint deck with dark gray background, one image per slide."""

import os
from pptx import Presentation
from pptx.util import Inches, Emu
from pptx.dml.color import RGBColor
from PIL import Image

IMAGE_DIR = "/Users/daniellefcourt/Desktop/arenaDownloader/arena_baroque-user-interface"
OUTPUT = "/Users/daniellefcourt/Desktop/arenaDownloader/baroque-user-interface-slides.pptx"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
BG_COLOR = RGBColor(0x33, 0x33, 0x33)  # dark gray (#333333)

SUPPORTED_EXT = {".jpg", ".jpeg", ".png", ".gif", ".bmp", ".tiff"}

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

# Blank layout
blank_layout = prs.slide_layouts[6]

files = sorted(os.listdir(IMAGE_DIR))
images = [f for f in files if os.path.splitext(f)[1].lower() in SUPPORTED_EXT]

print(f"Found {len(images)} images")

for fname in images:
    fpath = os.path.join(IMAGE_DIR, fname)

    # Get image dimensions
    try:
        with Image.open(fpath) as img:
            img_w, img_h = img.size
    except Exception as e:
        print(f"  Skipping {fname}: {e}")
        continue

    slide = prs.slides.add_slide(blank_layout)

    # Set background to dark gray
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

    # Scale image to fit slide while preserving aspect ratio
    slide_w_px = SLIDE_W
    slide_h_px = SLIDE_H

    scale_w = slide_w_px / Emu(int(img_w * 914400))  # 1 inch = 914400 EMU, 1px ≈ 1/96 inch but we use ratio
    scale_h = slide_h_px / Emu(int(img_h * 914400))

    # Simpler: compute based on aspect ratios
    img_aspect = img_w / img_h
    slide_aspect = SLIDE_W / SLIDE_H

    if img_aspect > slide_aspect:
        # Image is wider than slide — fit to width
        pic_w = SLIDE_W
        pic_h = int(SLIDE_W / img_aspect)
    else:
        # Image is taller than slide — fit to height
        pic_h = SLIDE_H
        pic_w = int(SLIDE_H * img_aspect)

    left = (SLIDE_W - pic_w) // 2
    top = (SLIDE_H - pic_h) // 2

    slide.shapes.add_picture(fpath, left, top, pic_w, pic_h)
    print(f"  Added: {fname}")

prs.save(OUTPUT)
print(f"\nSaved {len(images)} slides to {OUTPUT}")
